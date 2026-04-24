"""
Exportador PPT — gera apresentação PowerPoint executiva.

Requer: pip install python-pptx matplotlib

Uso:
    pptx_bytes = ExportadorPPT.gerar(kpis, df_mensal, df_dre, df_pareto, titulo)
    Path("apresentacao.pptx").write_bytes(pptx_bytes)
"""
from __future__ import annotations

import io
import logging
from typing import Any, Optional

import pandas as pd

logger = logging.getLogger(__name__)

_AZUL = "1F4E79"
_VERDE = "1E8449"
_VERMELHO = "C0392B"
_CINZA = "F2F2F2"


class ExportadorPPT:
    """Gera slides executivos com KPIs, gráfico e tabelas."""

    @staticmethod
    def gerar(
        kpis: dict[str, Any],
        df_mensal: Optional[pd.DataFrame] = None,
        df_dre: Optional[pd.DataFrame] = None,
        df_pareto: Optional[pd.DataFrame] = None,
        titulo: str = "Análise Financeira",
        nome_empresa: str = "",
    ) -> bytes:
        """Retorna bytes do arquivo .pptx."""
        try:
            from pptx import Presentation  # noqa: PLC0415
            from pptx.util import Inches, Pt
        except ImportError as exc:
            raise ImportError(
                "python-pptx não instalado. Execute: pip install python-pptx matplotlib"
            ) from exc

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        _slide_capa(prs, titulo, nome_empresa)
        _slide_kpis(prs, kpis)
        if df_mensal is not None and len(df_mensal) > 0:
            _slide_grafico(prs, df_mensal)
        if df_dre is not None and len(df_dre) > 0:
            _slide_dre(prs, df_dre)
        if df_pareto is not None and len(df_pareto) > 0:
            _slide_pareto(prs, df_pareto)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()


def _slide_capa(prs: Any, titulo: str, empresa: str) -> None:
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)

    tf = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf.text_frame.text = titulo
    p = tf.text_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    if empresa:
        tf2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(0.8))
        tf2.text_frame.text = empresa
        tf2.text_frame.paragraphs[0].font.size = Pt(22)
        tf2.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xCC, 0xCC, 0xFF)


def _slide_kpis(prs: Any, kpis: dict) -> None:
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _titulo_slide(slide, "Indicadores Financeiros")

    cards = [
        ("Receita Total", kpis.get("receita_total", 0), _VERDE),
        ("Despesa Total", kpis.get("despesa_total", 0), _VERMELHO),
        ("Resultado Líquido", kpis.get("resultado_liquido", 0), _AZUL),
        ("NFs Vendidas", kpis.get("nfs_receita", 0), _VERDE),
        ("NFs Recebidas", kpis.get("nfs_despesa", 0), _VERMELHO),
        ("Margem %", kpis.get("margem_pct", 0), _AZUL),
    ]

    col_w, row_h = Inches(3.8), Inches(1.8)
    x_start, y_start = Inches(0.3), Inches(1.2)
    for i, (label, valor, cor_hex) in enumerate(cards):
        col = i % 3
        row = i // 3
        x = x_start + col * (col_w + Inches(0.2))
        y = y_start + row * (row_h + Inches(0.15))
        shape = slide.shapes.add_shape(1, x, y, col_w, row_h)  # MSO_SHAPE_TYPE.RECTANGLE
        shape.fill.solid()
        r, g, b = int(cor_hex[0:2], 16), int(cor_hex[2:4], 16), int(cor_hex[4:6], 16)
        shape.fill.fore_color.rgb = RGBColor(r, g, b)
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        valor_str = f"R$ {valor:,.2f}".replace(",", "X").replace(".", ",").replace("X", ".") if isinstance(valor, float) else str(valor)
        tf.text = f"{label}\n{valor_str}"
        tf.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.bold = True
        if len(tf.paragraphs) > 1:
            tf.paragraphs[1].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            tf.paragraphs[1].font.size = Pt(20)


def _slide_grafico(prs: Any, df_mensal: pd.DataFrame) -> None:
    from pptx.util import Inches

    try:
        import matplotlib.pyplot as plt  # noqa: PLC0415
    except ImportError:
        return

    fig, ax = plt.subplots(figsize=(11, 4.5))
    periodos = df_mensal["Periodo"].astype(str).tolist()
    receitas = df_mensal["Receita_RS"].tolist()
    despesas = df_mensal["Despesa_RS"].tolist()
    x = range(len(periodos))
    ax.bar([i - 0.2 for i in x], receitas, width=0.4, label="Receitas", color="#1E8449")
    ax.bar([i + 0.2 for i in x], despesas, width=0.4, label="Despesas", color="#C0392B")
    ax.set_xticks(list(x))
    ax.set_xticklabels(periodos, rotation=45, ha="right", fontsize=9)
    ax.set_title("Receitas vs Despesas por Período")
    ax.legend()
    fig.tight_layout()

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150)
    plt.close(fig)
    buf.seek(0)

    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _titulo_slide(slide, "Receitas vs Despesas")
    slide.shapes.add_picture(buf, Inches(0.5), Inches(1.2), Inches(12), Inches(5.5))


def _slide_dre(prs: Any, df_dre: pd.DataFrame) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _titulo_slide(slide, "DRE Simplificada")
    _tabela_pptx(slide, df_dre.head(12))


def _slide_pareto(prs: Any, df_pareto: pd.DataFrame) -> None:
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _titulo_slide(slide, "Top Clientes / Fornecedores")
    _tabela_pptx(slide, df_pareto.head(10))


def _titulo_slide(slide: Any, texto: str) -> None:
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    tf = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.7))
    tf.text_frame.text = texto
    p = tf.text_frame.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)


def _tabela_pptx(slide: Any, df: pd.DataFrame) -> None:
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    rows, cols = len(df) + 1, len(df.columns)
    table = slide.shapes.add_table(rows, cols, Inches(0.3), Inches(1.1), Inches(12.5), Inches(5.8)).table

    for c, col_name in enumerate(df.columns):
        cell = table.cell(0, c)
        cell.text = str(col_name)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(11)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    for r, (_, row) in enumerate(df.iterrows(), start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(10)
